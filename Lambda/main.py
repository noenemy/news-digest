from bs4 import BeautifulSoup
import requests
import datetime as dt
from pptx import Presentation


def get_whatsnew_url_list(year, month):
    url_list = []

    base_url = "https://aws.amazon.com/about-aws/whats-new"
    url = f"{base_url}/{year}/{'{:02}'.format(month)}"

    whatsnew_list_page = requests.get(url)
    soup = BeautifulSoup(whatsnew_list_page.text, "html.parser")
    whatsnew_list = soup.select(selector="ul li h3 a")

    url_list = ["https:" + item.get("href") for item in whatsnew_list]

    return url_list


def read_whatsnew_article(url):

    whatsnew_page = requests.get(url)

    soup = BeautifulSoup(whatsnew_page.text, "html.parser")
    title = soup.select_one(selector="div h1 a").get_text()
    posted_date = soup.find(name="span", class_="date").get_text().strip()
    posted_date = dt.datetime.strptime(posted_date, "%b %d, %Y")

    content = []
    textboxes = soup.find_all("div", class_="aws-text-box")
    #print(type(textboxes), textboxes)
    for textbox in textboxes:
        text = textbox.find("p").text
        content.append(text)

    article = {
        "url": url,
        "posted_date": posted_date.isoformat(),
        "title": title,
        "content": content,
    }

    return article


def get_article_sample():
    article = {
        "url": "https://aws.amazon.com/about-aws/whats-new/2021/02/amazon-vpc-traffic-mirroring-supported-select-non-nitro-instance-types/",
        "posted_date": "2021-02-10T00:00:00",
        "title": " Amazon VPC Traffic Mirroring is now supported on select non-Nitro instance types",
        "content": '[<p>Amazon Virtual Private Cloud (Amazon VPC) Traffic Mirroring is now supported on additional select non-Nitro instance types. Amazon VPC Traffic Mirroring allows you to replicate the network traffic from EC2 instances within your VPC to security and monitoring appliances for use cases such as content inspection, threat monitoring, and troubleshooting.</p>, <p>Until now, customers could only enable VPC Traffic Mirroring on their Nitro-based EC2 instances. With this release, customers can now enable VPC Traffic Mirroring on additional instances types (complete list below) that use the Xen-based hypervisor. This enables you to now uniformly inspect network traffic on these additional EC2 instance types.<br/> </p>, <p>List of instance types that are now supported with this release are C4, D2, G3, G3s, H1, I3, M4, P2, P3, R4, X1 and X1e. VPC Traffic Mirroring is not supported on the T2, C3, R3 and I2 instance types and previous generation instances.<br/> </p>, <p>These additional instance types are supported in all 20 regions where VPC Traffic Mirroring is currently supported. To learn more about VPC Traffic Mirroring, please visit the VPC Traffic Mirroring <a href="https://docs.aws.amazon.com/vpc/latest/mirroring/what-is-traffic-mirroring.html" target="_blank">documentation</a>.  </p>]',
    }

    return article


def make_ppt(article, file_path):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(title_slide_layout)

    title = slide.placeholders[0]
    title.text = article["title"]

    body = slide.placeholders[1]
    tf = body.text_frame

    for paragraph in article["content"]:
        p = tf.add_paragraph()
        p.text = paragraph

    prs.save(file_path)

def Hyperlink( run_object, source_slide, destination_slide ):
    rId = source_slide.part.relate_to(destination_slide.part, RT.SLIDE)
    rPr = run_object._r.get_or_add_rPr()
    hlinkClick = rPr.add_hlinkClick(rId)
    hlinkClick.set('action', 'ppaction://hlinksldjump')

year = dt.datetime.now().year
month = dt.datetime.now().month

url_list = get_whatsnew_url_list(year, month)
article = read_whatsnew_article(url_list[0])

# article = get_article_sample()
# print(article)

make_ppt(article, "whatsnew_sample.pptx")