import json
from bs4 import BeautifulSoup, element
import requests
import datetime as dt
from pptx import Presentation
import boto3

def get_whatsnew_url_list(year, month):
    url_list = []

    base_url = "https://aws.amazon.com/about-aws/whats-new"
    url = f"{base_url}/{year}/{'{:02}'.format(month)}"

    whatsnew_list_page = requests.get(url)
    soup = BeautifulSoup(whatsnew_list_page.text, "html.parser")
    whatsnew_list = soup.select(selector="ul li h3 a")

    url_list = ["https:" + item.get("href") for item in whatsnew_list]

    return url_list

def read_whatsnew_article(url):

    whatsnew_page = requests.get(url)

    soup = BeautifulSoup(whatsnew_page.text, "html.parser")
    title = soup.select_one(selector="div h1 a").get_text()
    posted_date = soup.find(name="span", class_="date").get_text().strip()
    posted_date = dt.datetime.strptime(posted_date, "%b %d, %Y")

    content = []
    textboxes = soup.find_all("div", class_="aws-text-box")
    # print(type(textboxes), textboxes)
    for textbox in textboxes:

        paragraph = textbox.find("p")
        # HTML tag 단위로 분리해서 리스트로 구성한다.
        for child in paragraph.children:
            content.append(child)

    # print(content)

    article = {
        "url": url,
        "posted_date": posted_date.isoformat(),
        "title": title,
        "content": content,
    }

    return article

def get_aws_url(url):
    if url[0] == "/":
        url = "https://aws.amazon.com" + url
    return url

def get_article_sample():
    article = {
        "url": "https://aws.amazon.com/about-aws/whats-new/2021/02/amazon-vpc-traffic-mirroring-supported-select-non-nitro-instance-types/",
        "posted_date": "2021-02-10T00:00:00",
        "title": " Amazon VPC Traffic Mirroring is now supported on select non-Nitro instance types",
        "content": '[<p>Amazon Virtual Private Cloud (Amazon VPC) Traffic Mirroring is now supported on additional select non-Nitro instance types. Amazon VPC Traffic Mirroring allows you to replicate the network traffic from EC2 instances within your VPC to security and monitoring appliances for use cases such as content inspection, threat monitoring, and troubleshooting.</p>, <p>Until now, customers could only enable VPC Traffic Mirroring on their Nitro-based EC2 instances. With this release, customers can now enable VPC Traffic Mirroring on additional instances types (complete list below) that use the Xen-based hypervisor. This enables you to now uniformly inspect network traffic on these additional EC2 instance types.<br/> </p>, <p>List of instance types that are now supported with this release are C4, D2, G3, G3s, H1, I3, M4, P2, P3, R4, X1 and X1e. VPC Traffic Mirroring is not supported on the T2, C3, R3 and I2 instance types and previous generation instances.<br/> </p>, <p>These additional instance types are supported in all 20 regions where VPC Traffic Mirroring is currently supported. To learn more about VPC Traffic Mirroring, please visit the VPC Traffic Mirroring <a href="https://docs.aws.amazon.com/vpc/latest/mirroring/what-is-traffic-mirroring.html" target="_blank">documentation</a>.  </p>]',
    }

    return article

def make_slide(article, slide):
    title = slide.placeholders[0]
    title.text = article["title"]

    body = slide.placeholders[1]
    tf = body.text_frame
    p = tf.add_paragraph()

    for item in article["content"]:

        item_type = type(item)
        print(item_type)
        print(f"{item}")
        print("-------")

        if item.name == "a":
            print("add hyperlink")
            print("-------")
            run = p.add_run()
            run.text = item.text
            run.hyperlink.address = get_aws_url(item["href"])
        elif item.name == "br":
            print("add line break")
            print("-------")
            p.add_line_break()
            p.add_line_break()
        elif item.name == "b":
            run = p.add_run()
            run.text = item.text
            run.font.bold = True
        else:
            run = p.add_run()
            run.text = item

def make_ppt(url_list, file_path):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[1]

    max = len(url_list)
    if max > 5: # TODO : for testing only
        max = 5
    for index in range(0, max):
        article = read_whatsnew_article(url_list[index])

        slide = prs.slides.add_slide(title_slide_layout)
        make_slide(article, slide)

    tmp_file_path = f"/tmp/{file_path}"
    prs.save(tmp_file_path)

    s3_bucket_name = "soonki-shared-seoul" #TODO:
    s3_client = boto3.client("s3")
    s3_client.upload_file(tmp_file_path, s3_bucket_name, file_path)

    presigned_url = s3_client.generate_presigned_url(
        'get_object',
        Params={
            "Bucket": s3_bucket_name,
            "Key": file_path
        },
        HttpMethod="GET"
    )
    print(presigned_url)
    return presigned_url

def lambda_handler(event, context):
    """Sample pure Lambda function

    Parameters
    ----------
    event: dict, required
        API Gateway Lambda Proxy Input Format

        Event doc: https://docs.aws.amazon.com/apigateway/latest/developerguide/set-up-lambda-proxy-integrations.html#api-gateway-simple-proxy-for-lambda-input-format

    context: object, required
        Lambda Context runtime methods and attributes

        Context doc: https://docs.aws.amazon.com/lambda/latest/dg/python-context-object.html

    Returns
    ------
    API Gateway Lambda Proxy Output Format: dict

        Return doc: https://docs.aws.amazon.com/apigateway/latest/developerguide/set-up-lambda-proxy-integrations.html
    """

    # try:
    #     ip = requests.get("http://checkip.amazonaws.com/")
    # except requests.RequestException as e:
    #     # Send some context about this error to Lambda Logs
    #     print(e)

    #     raise e

    year = dt.datetime.now().year
    month = dt.datetime.now().month

    url_list = get_whatsnew_url_list(year, month)

    # article = get_article_sample()

    download_url = make_ppt(url_list, "whatsnew_sample.pptx")    

    return {
        "statusCode": 200,
        "body": json.dumps({
            "message": download_url,
            # "location": ip.text.replace("\n", "")
        }),
    }
